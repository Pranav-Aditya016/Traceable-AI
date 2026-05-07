from pptx import Presentation

pptx_path = r'C:\Pranav Aditya\Traceable AI\Pranav Aditya PPT 2nd review.pptx'
prs = Presentation(pptx_path)

print(f'Total slides: {len(prs.slides)}')
print('='*80)

for i, slide in enumerate(prs.slides, 1):
    layout_name = slide.slide_layout.name if slide.slide_layout else 'Unknown'
    print(f'\n### SLIDE {i} (Layout: {layout_name})')
    print('-'*60)
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    print(f'  {text}')
        if shape.has_table:
            table = shape.table
            rows = len(table.rows)
            cols = len(table.columns)
            print(f'  [TABLE {rows} rows x {cols} cols]')
            for row in table.rows:
                row_text = ' | '.join(cell.text.strip() for cell in row.cells)
                print(f'    {row_text}')
        if shape.shape_type == 13:
            print(f'  [IMAGE: {shape.name}]')
    print()
